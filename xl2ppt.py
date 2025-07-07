import openpyxl
from pptx import Presentation
from pptx.util import Inches

# --- CONFIGURATION ---
excel_file = "Table_Template2.xlsx"
sheet_name = "Sheet2"  # Worksheet 2 (name or index)
data_range = "A1:D5"   # Change as needed
pptx_file = "slidetemplate.pptx"
output_pptx = "updated_presentation.pptx"
slide_index = 0  # 0-based index (Slide 1)

# --- Load Excel Data ---
wb = openpyxl.load_workbook(excel_file, data_only=True)
ws = wb[sheet_name]

#If your Excel sheet name is unknown, you can access by index:
#ws = wb.worksheets[1] (worksheet 2, since index is 0-based).

#If you need to paste into a specific table (e.g. second table on slide), modify the for shape in slide.shapes logic.
#
#


min_col = openpyxl.utils.cell.column_index_from_string(data_range.split(":")[0][0])
min_row = int(data_range.split(":")[0][1:])
max_col = openpyxl.utils.cell.column_index_from_string(data_range.split(":")[1][0])
max_row = int(data_range.split(":")[1][1:])

data = []
for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=True):
    data.append(row)

# --- Load PowerPoint and Find Table on Slide 1 ---
prs = Presentation(pptx_file)
slide = prs.slides[slide_index]

# Find the first table on the slide
table = None
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table is None:
    raise ValueError("No table found on the first slide.")

# --- Copy Data into Table ---
for i, row in enumerate(data):
    for j, val in enumerate(row):
        if i < len(table.rows) and j < len(table.columns):
            table.cell(i, j).text = str(val) if val is not None else ""

# --- Save Updated Presentation ---
prs.save(output_pptx)
print(f"Data copied successfully to {output_pptx}")
